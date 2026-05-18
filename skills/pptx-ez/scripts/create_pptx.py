#!/usr/bin/env python3
"""
create_pptx.py — Create a .pptx presentation from a JSON spec or a built-in demo.

Usage:
  python create_pptx.py OUTPUT [--spec SPEC_JSON]

Options:
  --spec PATH   JSON spec file (see SKILL.md for format)

Examples:
  python create_pptx.py demo.pptx
  python create_pptx.py report.pptx --spec spec.json
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Cm, Pt
except ImportError:
    print('Error: python-pptx not installed. Run: pip install python-pptx', file=sys.stderr)
    sys.exit(1)

# ── Constants ─────────────────────────────────────────────────────────────────

CHART_TYPE_MAP: dict[str, int] = {
    'bar':             XL_CHART_TYPE.BAR_CLUSTERED,
    'bar_stacked':     XL_CHART_TYPE.BAR_STACKED,
    'column':          XL_CHART_TYPE.COLUMN_CLUSTERED,
    'column_stacked':  XL_CHART_TYPE.COLUMN_STACKED,
    'line':            XL_CHART_TYPE.LINE,
    'line_markers':    XL_CHART_TYPE.LINE_MARKERS,
    'pie':             XL_CHART_TYPE.PIE,
    'doughnut':        XL_CHART_TYPE.DOUGHNUT,
    'area':            XL_CHART_TYPE.AREA,
    'scatter':         XL_CHART_TYPE.XY_SCATTER,
}

ALIGN_MAP: dict[str, PP_ALIGN] = {
    'left':    PP_ALIGN.LEFT,
    'center':  PP_ALIGN.CENTER,
    'right':   PP_ALIGN.RIGHT,
    'justify': PP_ALIGN.JUSTIFY,
}


# ── Helpers ───────────────────────────────────────────────────────────────────

def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _layout(prs: Presentation, name: str):
    name_map = {
        'title':       'Title Slide',
        'content':     'Title and Content',
        'two_content': 'Two Content',
        'title_only':  'Title Only',
        'blank':       'Blank',
    }
    target = name_map.get(name, name)
    for layout in prs.slide_layouts:
        if layout.name == target:
            return layout
    # fallback: index-based
    fallbacks = {'title': 0, 'content': 1, 'two_content': 3, 'title_only': 5, 'blank': 6}
    return prs.slide_layouts[fallbacks.get(name, 6)]


def _set_notes(slide, text: str) -> None:
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = text


def _apply_run_format(run, spec: dict) -> None:
    font = run.font
    if spec.get('bold'):
        font.bold = True
    if spec.get('italic'):
        font.italic = True
    if spec.get('underline'):
        font.underline = True
    if 'font_size' in spec:
        font.size = Pt(spec['font_size'])
    if 'font_color' in spec:
        font.color.rgb = _rgb(spec['font_color'])
    if 'font_name' in spec:
        font.name = spec['font_name']


def _add_textbox(slide, spec: dict) -> None:
    left   = Cm(spec.get('left_cm', 1))
    top    = Cm(spec.get('top_cm', 1))
    width  = Cm(spec.get('width_cm', 20))
    height = Cm(spec.get('height_cm', 2))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True

    text  = spec.get('text', '')
    lines = text.split('\n') if '\n' in text else [text]
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        run = para.add_run()
        run.text = line
        _apply_run_format(run, spec)
        align = spec.get('align', 'left')
        para.alignment = ALIGN_MAP.get(align, PP_ALIGN.LEFT)


def _add_image(slide, spec: dict) -> None:
    path   = spec['path']
    left   = Cm(spec.get('left_cm', 2))
    top    = Cm(spec.get('top_cm', 2))
    width  = Cm(spec['width_cm']) if 'width_cm' in spec else None
    height = Cm(spec['height_cm']) if 'height_cm' in spec else None
    slide.shapes.add_picture(path, left, top, width=width, height=height)


def _add_table(slide, spec: dict) -> None:
    rows_data = spec.get('rows', [])
    if not rows_data:
        return
    nrows = len(rows_data)
    ncols = max(len(r) for r in rows_data)

    left   = Cm(spec.get('left_cm', 2))
    top    = Cm(spec.get('top_cm', 5))
    width  = Cm(spec.get('width_cm', 22))
    height = Cm(spec.get('height_cm', max(2, nrows * 1.2)))

    tbl_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl       = tbl_shape.table

    header_bg    = spec.get('header_bg', '#4472C4')
    header_fg    = spec.get('header_font_color', '#FFFFFF')
    alt_bg       = spec.get('alt_row_bg', None)

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            paras = cell.text_frame.paragraphs
            if paras:
                runs = paras[0].runs
                if not runs:
                    run = paras[0].add_run()
                    run.text = str(val)
                    runs = [run]
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(header_bg)
                    for run in runs:
                        run.font.bold      = True
                        run.font.color.rgb = _rgb(header_fg)
                elif alt_bg and r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(alt_bg)


def _add_chart(slide, spec: dict) -> None:
    chart_type_key = spec.get('chart_type', 'bar')
    xl_type        = CHART_TYPE_MAP.get(chart_type_key, XL_CHART_TYPE.BAR_CLUSTERED)

    left   = Cm(spec.get('left_cm', 3))
    top    = Cm(spec.get('top_cm', 5))
    width  = Cm(spec.get('width_cm', 24))
    height = Cm(spec.get('height_cm', 12))

    chart_data = CategoryChartData()
    chart_data.categories = spec.get('categories', [])
    for series in spec.get('series', []):
        chart_data.add_series(series['name'], tuple(series['values']))

    frame = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = frame.chart
    chart.has_legend = spec.get('legend', True)
    if spec.get('title'):
        chart.has_title = True
        chart.chart_title.text_frame.text = spec['title']


def _add_shape_to_slide(slide, spec: dict) -> None:
    kind = spec.get('type', '')
    if kind == 'textbox':
        _add_textbox(slide, spec)
    elif kind == 'image':
        _add_image(slide, spec)
    elif kind == 'table':
        _add_table(slide, spec)
    elif kind == 'chart':
        _add_chart(slide, spec)


# ── Slide builders ────────────────────────────────────────────────────────────

def _build_title_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(_layout(prs, 'title'))
    if spec.get('title'):
        slide.shapes.title.text = spec['title']
    if spec.get('subtitle') and len(slide.placeholders) > 1:
        slide.placeholders[1].text = spec['subtitle']
    _set_notes(slide, spec.get('notes', ''))


def _build_content_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(_layout(prs, 'content'))
    if spec.get('title'):
        slide.shapes.title.text = spec['title']
    content = spec.get('content', [])
    if content and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(content):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text  = bullet
            para.level = spec.get('level', 0)
    _set_notes(slide, spec.get('notes', ''))


def _build_two_content_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(_layout(prs, 'two_content'))
    if spec.get('title'):
        slide.shapes.title.text = spec['title']
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    def fill_ph(idx, lines):
        if idx not in placeholders:
            return
        tf = placeholders[idx].text_frame
        tf.clear()
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line

    fill_ph(1, spec.get('left', []))
    fill_ph(2, spec.get('right', []))
    _set_notes(slide, spec.get('notes', ''))


def _build_title_only_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(_layout(prs, 'title_only'))
    if spec.get('title'):
        slide.shapes.title.text = spec['title']
    for shape_spec in spec.get('shapes', []):
        _add_shape_to_slide(slide, shape_spec)
    _set_notes(slide, spec.get('notes', ''))


def _build_blank_slide(prs: Presentation, spec: dict) -> None:
    slide = prs.slides.add_slide(_layout(prs, 'blank'))
    for shape_spec in spec.get('shapes', []):
        _add_shape_to_slide(slide, shape_spec)
    _set_notes(slide, spec.get('notes', ''))


# ── Entry point ───────────────────────────────────────────────────────────────

def build_presentation(spec: dict) -> Presentation:
    prs = Presentation()
    if spec.get('widescreen', True):
        prs.slide_width  = Cm(33.867)
        prs.slide_height = Cm(19.05)

    props = prs.core_properties
    if spec.get('title'):
        props.title = spec['title']
    if spec.get('author'):
        props.author = spec['author']

    builders = {
        'title':       _build_title_slide,
        'content':     _build_content_slide,
        'two_content': _build_two_content_slide,
        'title_only':  _build_title_only_slide,
        'blank':       _build_blank_slide,
    }

    for slide_spec in spec.get('slides', []):
        layout = slide_spec.get('layout', 'content')
        builder = builders.get(layout, _build_blank_slide)
        builder(prs, slide_spec)

    return prs


def _demo_spec() -> dict:
    return {
        'title': 'Demo Presentation',
        'author': 'Claude',
        'widescreen': True,
        'slides': [
            {
                'layout': 'title',
                'title': 'pptx-ez Demo',
                'subtitle': 'A python-pptx powered skill',
            },
            {
                'layout': 'content',
                'title': 'Key Features',
                'content': [
                    'Create from JSON spec',
                    'Multiple slide layouts',
                    'Text, images, tables, charts',
                    'Speaker notes',
                    'LibreOffice PDF export',
                ],
                'notes': 'Walk through each point slowly.',
            },
            {
                'layout': 'blank',
                'shapes': [
                    {
                        'type': 'textbox',
                        'text': 'Sample Data Table',
                        'left_cm': 2, 'top_cm': 1, 'width_cm': 28, 'height_cm': 2,
                        'bold': True, 'font_size': 28, 'align': 'center',
                    },
                    {
                        'type': 'table',
                        'rows': [
                            ['Region', 'Q1', 'Q2', 'Q3'],
                            ['North',  '120', '135', '150'],
                            ['South',   '95', '110', '125'],
                            ['East',   '105', '120', '138'],
                        ],
                        'left_cm': 2, 'top_cm': 4, 'width_cm': 28, 'height_cm': 7,
                    },
                ],
            },
            {
                'layout': 'blank',
                'shapes': [
                    {
                        'type': 'chart',
                        'chart_type': 'column',
                        'title': 'Quarterly Revenue',
                        'categories': ['Q1', 'Q2', 'Q3', 'Q4'],
                        'series': [
                            {'name': 'North', 'values': [120, 135, 150, 162]},
                            {'name': 'South', 'values': [ 95, 110, 125, 140]},
                        ],
                        'left_cm': 3, 'top_cm': 2, 'width_cm': 27, 'height_cm': 15,
                    },
                ],
            },
        ],
    }


def main() -> None:
    ap = argparse.ArgumentParser(
        description=__doc__,
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    ap.add_argument('output', help='Destination .pptx')
    ap.add_argument('--spec', default=None, help='JSON spec file')
    args = ap.parse_args()

    if args.spec:
        with open(args.spec, encoding='utf-8') as f:
            spec = json.load(f)
    else:
        spec = _demo_spec()

    prs = build_presentation(spec)
    prs.save(args.output)
    print(f'Saved: {args.output}')


if __name__ == '__main__':
    main()
