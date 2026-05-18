#!/usr/bin/env python3
"""
edit_pptx.py — Edit an existing .pptx file via python-pptx and direct OOXML manipulation.

Usage:
  python edit_pptx.py INPUT OUTPUT [options]

Options:
  --replace OLD NEW              Plain text search-replace across all slides
  --add-note SLIDE_NUM TEXT      Add/replace speaker note on slide N (1-based)
  --delete-slide SLIDE_NUM       Delete slide N (1-based)
  --add-slide LAYOUT             Append a new slide (layout: title|content|blank)
  --add-textbox SLIDE_NUM TEXT   Add a text box to slide N
    --left INCHES                  Left offset in inches (default 0.5)
    --top  INCHES                  Top offset in inches (default 0.5)
    --width INCHES                 Width in inches (default 8)
    --height INCHES                Height in inches (default 1)
    --font-size PT                 Font size in points (default 18)
    --bold                         Bold text
  --duplicate-slide SLIDE_NUM    Append a copy of slide N

Examples:
  python edit_pptx.py in.pptx out.pptx --replace "DRAFT" "FINAL"
  python edit_pptx.py in.pptx out.pptx --add-note 2 "Remember to demo the chart."
  python edit_pptx.py in.pptx out.pptx --delete-slide 3
  python edit_pptx.py in.pptx out.pptx --add-slide blank
  python edit_pptx.py in.pptx out.pptx --add-textbox 1 "Confidential" --top 6.5 --font-size 10
  python edit_pptx.py in.pptx out.pptx --duplicate-slide 1
"""
from __future__ import annotations

import argparse
import copy
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from lxml import etree
except ImportError as e:
    print(f'Error: {e}. Run: pip install python-pptx lxml', file=sys.stderr)
    sys.exit(1)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _layout(prs: Presentation, name: str):
    name_map = {
        'title':      'Title Slide',
        'content':    'Title and Content',
        'title_only': 'Title Only',
        'blank':      'Blank',
    }
    target = name_map.get(name, name)
    for layout in prs.slide_layouts:
        if layout.name == target:
            return layout
    fallbacks = {'title': 0, 'content': 1, 'title_only': 5, 'blank': 6}
    return prs.slide_layouts[fallbacks.get(name, 6)]


def _iter_runs(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        yield run
        notes = slide.notes_slide
        for para in notes.notes_text_frame.paragraphs:
            for run in para.runs:
                yield run


# ── Operations ────────────────────────────────────────────────────────────────

def replace(prs: Presentation, old: str, new: str) -> int:
    """Replace all occurrences of old with new across every run in every slide."""
    count = 0
    for run in _iter_runs(prs):
        if old in run.text:
            run.text = run.text.replace(old, new)
            count += 1
    return count


def add_note(prs: Presentation, slide_num: int, text: str) -> None:
    slide = prs.slides[slide_num - 1]
    slide.notes_slide.notes_text_frame.text = text


def delete_slide(prs: Presentation, slide_num: int) -> None:
    idx = slide_num - 1
    xml_slides = prs.slides._sldIdLst
    slide = prs.slides[idx]
    rId = prs.slides._sldIdLst[idx].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    xml_slides.remove(xml_slides[idx])


def add_slide(prs: Presentation, layout_name: str) -> None:
    prs.slides.add_slide(_layout(prs, layout_name))


def add_textbox(
    prs: Presentation,
    slide_num: int,
    text: str,
    left: float = 0.5,
    top: float  = 0.5,
    width: float = 8.0,
    height: float = 1.0,
    font_size: float = 18.0,
    bold: bool = False,
) -> None:
    slide = prs.slides[slide_num - 1]
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf   = txBox.text_frame
    para = tf.paragraphs[0]
    run  = para.add_run()
    run.text       = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold


def duplicate_slide(prs: Presentation, slide_num: int) -> None:
    """Append a copy of slide N to the end of the presentation."""
    template = prs.slides[slide_num - 1]
    try:
        layout = template.slide_layout
    except (KeyError, IndexError):
        raise ValueError(
            "Slide has no slideLayout relationship — cannot duplicate via python-pptx. "
            "Use scripts/add_slide.py on an unpacked presentation instead."
        )
    blank    = prs.slides.add_slide(layout)

    # Copy the XML tree of the source slide's spTree into the new slide
    src_spTree = template.shapes._spTree
    dst_spTree = blank.shapes._spTree

    # Remove default placeholders from blank slide
    for el in list(dst_spTree):
        dst_spTree.remove(el)
    for el in src_spTree:
        dst_spTree.append(copy.deepcopy(el))

    # Copy notes if present
    src_notes = template.notes_slide.notes_text_frame.text
    if src_notes:
        blank.notes_slide.notes_text_frame.text = src_notes


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> None:
    ap = argparse.ArgumentParser(
        description=__doc__,
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    ap.add_argument('input',  help='Source .pptx')
    ap.add_argument('output', help='Destination .pptx')
    ap.add_argument('--replace',
                    nargs=2, metavar=('OLD', 'NEW'),
                    action='append', default=[])
    ap.add_argument('--add-note',
                    nargs=2, metavar=('SLIDE_NUM', 'TEXT'),
                    action='append', default=[], dest='add_note')
    ap.add_argument('--delete-slide',
                    nargs=1, metavar='SLIDE_NUM',
                    action='append', default=[], dest='delete_slide')
    ap.add_argument('--add-slide',
                    nargs=1, metavar='LAYOUT',
                    action='append', default=[], dest='add_slide')
    ap.add_argument('--add-textbox',
                    nargs=2, metavar=('SLIDE_NUM', 'TEXT'),
                    action='append', default=[], dest='add_textbox')
    ap.add_argument('--left',      type=float, default=0.5)
    ap.add_argument('--top',       type=float, default=0.5)
    ap.add_argument('--width',     type=float, default=8.0)
    ap.add_argument('--height',    type=float, default=1.0)
    ap.add_argument('--font-size', type=float, default=18.0, dest='font_size')
    ap.add_argument('--bold',      action='store_true')
    ap.add_argument('--duplicate-slide',
                    nargs=1, metavar='SLIDE_NUM',
                    action='append', default=[], dest='duplicate_slide')
    args = ap.parse_args()

    prs = Presentation(args.input)

    for old, new in args.replace:
        n = replace(prs, old, new)
        print(f'Replaced {n} run(s): {old!r} → {new!r}')

    for num_str, text in args.add_note:
        add_note(prs, int(num_str), text)

    # Process deletions in reverse order to keep indices stable
    for (num_str,) in sorted(args.delete_slide, key=lambda x: int(x[0]), reverse=True):
        delete_slide(prs, int(num_str))

    for (layout,) in args.add_slide:
        add_slide(prs, layout)

    for num_str, text in args.add_textbox:
        add_textbox(
            prs, int(num_str), text,
            left=args.left, top=args.top,
            width=args.width, height=args.height,
            font_size=args.font_size, bold=args.bold,
        )

    for (num_str,) in args.duplicate_slide:
        duplicate_slide(prs, int(num_str))

    prs.save(args.output)
    print(f'Saved: {args.output}')


if __name__ == '__main__':
    main()
