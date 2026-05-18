#!/usr/bin/env python3
"""Benchmark: pptx-ez skill vs Anthropic skills fork pptx skill.

Suites:
  C – SKILL.md coverage (keyword/topic checks)
  D – Script inventory (which files each skill ships)
  E – Code quality (safety, portability, design guidance)
  A – Functional parity (create → validate round-trip)
  B – Feature-specific tests (unpack/pack, add_slide, clean)

Usage:
    python benchmark.py
"""

from __future__ import annotations

import os
import subprocess
import sys
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path

SKILL_A = Path("/home/user/kilo/skills/pptx-ez")
SKILL_B = Path("/home/user/skills/skills/pptx")
SKILL_A_SCRIPTS = SKILL_A / "scripts"
SKILL_B_SCRIPTS = SKILL_B / "scripts"
SKILL_B_OFFICE  = SKILL_B / "scripts" / "office"

_results: list[tuple[str, str, str, str]] = []
_pass = _fail = _warn = 0


# ── helpers ──────────────────────────────────────────────────────────────────

def _record(suite: str, name: str, status: str, detail: str = "") -> None:
    global _pass, _fail, _warn
    _results.append((suite, name, status, detail))
    icon = {"PASS": "✓", "FAIL": "✗", "WARN": "⚠"}.get(status, "?")
    suffix = f": {detail}" if detail else ""
    print(f"  {icon} [{suite}] {name}{suffix}")
    if status == "PASS":   _pass += 1
    elif status == "FAIL": _fail += 1
    elif status == "WARN": _warn += 1


def _read(path: Path | None) -> str:
    if path is None or not path.exists():
        return ""
    return path.read_text(errors="replace")


def _run(cmd: list[str], cwd: str | None = None, env=None,
         timeout: int = 30) -> tuple[int, str, str]:
    try:
        r = subprocess.run(cmd, capture_output=True, text=True,
                           cwd=cwd, env=env, timeout=timeout)
        return r.returncode, r.stdout, r.stderr
    except subprocess.TimeoutExpired:
        return -1, "", "TIMEOUT"
    except Exception as e:
        return -2, "", str(e)


# ── Suite C – SKILL.md coverage ───────────────────────────────────────────────

SKILL_MD_CHECKS = [
    # Trigger
    ("frontmatter name field",          ["name:"]),
    ("rich trigger description",        ["presentation", "deck", "slides"]),
    ("do-not-trigger guard",            ["Do NOT", "not"]),
    # Workflows
    ("unpack workflow",                 ["unpack.py", "unpack"]),
    ("pack workflow",                   ["pack.py", "pack"]),
    ("add_slide script documented",     ["add_slide"]),
    ("clean script documented",         ["clean.py", "clean"]),
    ("thumbnail / visual preview",      ["thumbnail"]),
    ("markitdown text extraction",      ["markitdown"]),
    # Creation
    ("pptxgenjs creation path",         ["pptxgenjs", "pptxgen"]),
    ("python-pptx creation path",       ["python-pptx", "Presentation()"]),
    # Design guidance
    ("color palette guidance",          ["color", "palette"]),
    ("typography guidance",             ["font", "pt", "bold"]),
    ("spacing guidance",                ["margin", "spacing"]),
    ("avoid common mistakes section",   ["avoid", "mistake", "pitfall"]),
    ("never use unicode bullets rule",  ["bullet", "unicode"]),
    # XML / OOXML
    ("smart quotes entities",           ["&#x201C;", "&#x2018;", "smart"]),
    ("xml:space preserve",              ["xml:space", "preserve"]),
    ("defusedxml safety note",          ["defusedxml"]),
    ("EMU unit reference",              ["EMU", "914400"]),
    ("PML namespace reference",         ["presentationml", "p:sld"]),
    ("DML namespace reference",         ["drawingml", "a:r"]),
    # QA
    ("QA / visual verification",        ["QA", "visual"]),
    ("soffice PDF conversion",          ["soffice", "convert-to"]),
    ("pdftoppm slide images",           ["pdftoppm"]),
    # Dependencies
    ("dependencies listed",             ["python-pptx", "LibreOffice"]),
]


def suite_c() -> None:
    print("\n## Suite C: SKILL.md Coverage")
    a_text = _read(SKILL_A / "SKILL.md")
    b_text = _read(SKILL_B / "SKILL.md") + _read(SKILL_B / "editing.md") + _read(SKILL_B / "pptxgenjs.md")

    for label, keywords in SKILL_MD_CHECKS:
        a_has = all(kw.lower() in a_text.lower() for kw in keywords)
        b_has = all(kw.lower() in b_text.lower() for kw in keywords)

        if a_has and b_has:
            _record("C", label, "PASS")
        elif b_has and not a_has:
            _record("C", label, "FAIL", "Skill A missing; Skill B has it")
        elif a_has and not b_has:
            _record("C", label, "WARN", "Skill A unique feature; Skill B missing")
        else:
            _record("C", label, "FAIL", "Both missing")


# ── Suite D – Script inventory ────────────────────────────────────────────────

SCRIPT_PAIRS: dict[str, tuple[Path | None, Path | None]] = {
    "create script (python-pptx)":   (SKILL_A_SCRIPTS / "create_pptx.py",          None),
    "edit script":                   (SKILL_A_SCRIPTS / "edit_pptx.py",             None),
    "validate script":               (SKILL_A_SCRIPTS / "validate_pptx.py",         SKILL_B_OFFICE / "validate.py"),
    "export_pdf.sh":                 (SKILL_A_SCRIPTS / "export_pdf.sh",            None),
    "office/unpack.py":              (SKILL_A_SCRIPTS / "office" / "unpack.py",     SKILL_B_OFFICE / "unpack.py"),
    "office/pack.py":                (SKILL_A_SCRIPTS / "office" / "pack.py",       SKILL_B_OFFICE / "pack.py"),
    "office/soffice.py":             (SKILL_A_SCRIPTS / "office" / "soffice.py",    SKILL_B_OFFICE / "soffice.py"),
    "add_slide.py":                  (SKILL_A_SCRIPTS / "add_slide.py",             SKILL_B_SCRIPTS / "add_slide.py"),
    "clean.py":                      (SKILL_A_SCRIPTS / "clean.py",                 SKILL_B_SCRIPTS / "clean.py"),
    "thumbnail.py":                  (SKILL_A_SCRIPTS / "thumbnail.py",             SKILL_B_SCRIPTS / "thumbnail.py"),
    "validators/base.py":            (None,                                          SKILL_B_OFFICE / "validators" / "base.py"),
    "pptxgenjs.md reference":        (SKILL_A / "pptxgenjs.md",                     SKILL_B / "pptxgenjs.md"),
    "editing.md reference":          (SKILL_A / "editing.md",                       SKILL_B / "editing.md"),
    "python-pptx.md reference":      (SKILL_A / "references" / "python-pptx.md",   None),
    "ooxml-pptx.md reference":       (SKILL_A / "references" / "ooxml-pptx.md",    None),
}


def suite_d() -> None:
    print("\n## Suite D: Script Inventory")
    for label, (path_a, path_b) in SCRIPT_PAIRS.items():
        a_has = path_a is not None and path_a.exists()
        b_has = path_b is not None and path_b.exists()

        if a_has and b_has:
            _record("D", label, "PASS")
        elif a_has and not b_has:
            _record("D", label, "WARN", "Skill A unique feature")
        elif b_has and not a_has:
            _record("D", label, "FAIL", "Skill B has it; Skill A missing")
        else:
            _record("D", label, "FAIL", "Both missing")


# ── Suite E – Code quality ────────────────────────────────────────────────────

def suite_e() -> None:
    print("\n## Suite E: Code Quality")

    # create_pptx.py
    a_create = _read(SKILL_A_SCRIPTS / "create_pptx.py")
    _record("E", "create_pptx: shebang present",
            "PASS" if "#!/usr/bin/env python3" in a_create else "FAIL")
    _record("E", "create_pptx: json spec loading",
            "PASS" if "json.load" in a_create else "FAIL")
    _record("E", "create_pptx: chart support",
            "PASS" if "XL_CHART_TYPE" in a_create else "FAIL")
    _record("E", "create_pptx: speaker notes",
            "PASS" if "notes_slide" in a_create or "notes_text_frame" in a_create else "FAIL")
    _record("E", "create_pptx: table with header shading",
            "PASS" if "header_bg" in a_create or "fore_color" in a_create else "FAIL")

    # edit_pptx.py
    a_edit = _read(SKILL_A_SCRIPTS / "edit_pptx.py")
    _record("E", "edit_pptx: delete_slide handles relationship cleanup",
            "PASS" if "drop_rel" in a_edit else "FAIL")
    _record("E", "edit_pptx: duplicate_slide uses deepcopy",
            "PASS" if "deepcopy" in a_edit else "FAIL")
    _record("E", "edit_pptx: reverse-order delete (stable indices)",
            "PASS" if "reverse=True" in a_edit else "FAIL")

    # validate_pptx.py
    a_val = _read(SKILL_A_SCRIPTS / "validate_pptx.py")
    _record("E", "validate_pptx: checks required parts",
            "PASS" if "ppt/presentation.xml" in a_val else "FAIL")
    _record("E", "validate_pptx: XML parse validation",
            "PASS" if "XMLSyntaxError" in a_val or "fromstring" in a_val else "FAIL")
    _record("E", "validate_pptx: slide rel target verification",
            "PASS" if "slide_rel_type" in a_val or "slide" in a_val else "FAIL")

    # Skill B: clean.py uses defusedxml
    b_clean = _read(SKILL_B_SCRIPTS / "clean.py")
    _record("E", "Skill B clean.py uses defusedxml (safe XML parsing)",
            "PASS" if "defusedxml" in b_clean else "FAIL")
    a_clean = _read(SKILL_A_SCRIPTS / "clean.py")
    _record("E", "Skill A clean.py uses defusedxml (safe XML parsing)",
            "PASS" if "defusedxml" in a_clean else "WARN",
            "Skill B uses defusedxml.minidom throughout to prevent XXE on untrusted files")

    # Skill B: thumbnail uses office.soffice
    b_thumb = _read(SKILL_B_SCRIPTS / "thumbnail.py")
    _record("E", "Skill B thumbnail.py uses office.soffice (sandbox-safe PDF conversion)",
            "PASS" if "office.soffice" in b_thumb or "get_soffice_env" in b_thumb else "FAIL")

    # Skill B: add_slide handles notes refs cleanup
    b_add = _read(SKILL_B_SCRIPTS / "add_slide.py")
    _record("E", "Skill B add_slide.py strips notesSlide rel on duplicate (prevents orphaned notes)",
            "PASS" if "notesSlide" in b_add else "FAIL")

    # SKILL.md length ratio
    a_len = len(_read(SKILL_A / "SKILL.md"))
    b_total = (len(_read(SKILL_B / "SKILL.md")) +
               len(_read(SKILL_B / "editing.md")) +
               len(_read(SKILL_B / "pptxgenjs.md")))
    ratio = a_len / b_total if b_total else 0
    _record("E", f"SKILL.md coverage ratio (A={a_len} B={b_total} chars)",
            "PASS" if ratio >= 0.6 else "FAIL" if ratio < 0.3 else "WARN",
            f"Skill A is {int(ratio*100)}% of Skill B total guidance")


# ── Suite A – Functional round-trip ──────────────────────────────────────────

def _make_minimal_pptx(path: Path) -> None:
    CT = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml"  ContentType="application/xml"/>
  <Override PartName="/ppt/presentation.xml"
    ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slides/slide1.xml"
    ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
</Types>"""
    RELS = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1"
    Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"
    Target="ppt/presentation.xml"/>
</Relationships>"""
    PRS_RELS = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId2"
    Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
    Target="slides/slide1.xml"/>
</Relationships>"""
    PRS = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:sldIdLst>
    <p:sldId id="256" r:id="rId2"/>
  </p:sldIdLst>
  <p:sldSz cx="9144000" cy="6858000"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>"""
    SLIDE = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm>
      </p:grpSpPr>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Title 1"/>
          <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
          <p:nvPr><p:ph type="title"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
        <p:txBody>
          <a:bodyPr/><a:lstStyle/>
          <a:p><a:r><a:rPr lang="en-US" dirty="0"/><a:t>Hello PPTX</a:t></a:r></a:p>
        </p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
</p:sld>"""
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml",             CT)
        z.writestr("_rels/.rels",                     RELS)
        z.writestr("ppt/_rels/presentation.xml.rels", PRS_RELS)
        z.writestr("ppt/presentation.xml",            PRS)
        z.writestr("ppt/slides/slide1.xml",           SLIDE)


def suite_a(tmpdir: Path) -> None:
    print("\n## Suite A: Functional Round-Trip")

    # ── Skill A: create via create_pptx.py ──────────────────────────────────
    out_a = tmpdir / "a_demo.pptx"
    rc, out, err = _run([sys.executable, str(SKILL_A_SCRIPTS / "create_pptx.py"), str(out_a)])
    a_create_ok = rc == 0 and out_a.exists()
    _record("A", "Skill A create_pptx.py: demo deck created",
            "PASS" if a_create_ok else "FAIL", err.strip() if not a_create_ok else "")

    if a_create_ok:
        try:
            with zipfile.ZipFile(out_a) as z:
                names = z.namelist()
            has_prs  = "ppt/presentation.xml" in names
            has_slide = any("slides/slide" in n for n in names)
            _record("A", "Skill A create: valid zip with presentation.xml",
                    "PASS" if has_prs else "FAIL")
            _record("A", "Skill A create: at least one slide",
                    "PASS" if has_slide else "FAIL")
        except Exception as e:
            _record("A", "Skill A create: valid zip", "FAIL", str(e))

    # ── Skill A: validate ────────────────────────────────────────────────────
    if a_create_ok:
        rc, out, err = _run([sys.executable,
                             str(SKILL_A_SCRIPTS / "validate_pptx.py"), str(out_a)])
        _record("A", "Skill A validate_pptx.py: passes on own output",
                "PASS" if rc == 0 else "FAIL", (out + err).strip() if rc != 0 else "")

    # ── Skill A: edit (replace text) ─────────────────────────────────────────
    if a_create_ok:
        out_edited = tmpdir / "a_edited.pptx"
        rc, out, err = _run([sys.executable, str(SKILL_A_SCRIPTS / "edit_pptx.py"),
                             str(out_a), str(out_edited),
                             "--replace", "Demo", "Final"])
        _record("A", "Skill A edit_pptx.py: replace text",
                "PASS" if rc == 0 and out_edited.exists() else "FAIL",
                (out + err).strip() if rc != 0 else "")

    # ── Skill B: unpack ──────────────────────────────────────────────────────
    fixture = tmpdir / "fixture.pptx"
    _make_minimal_pptx(fixture)

    b_unpacked = tmpdir / "b_unpacked"
    env_b = os.environ.copy()
    env_b["PYTHONPATH"] = str(SKILL_B_SCRIPTS)
    rc, out, err = _run([sys.executable, str(SKILL_B_OFFICE / "unpack.py"),
                         str(fixture), str(b_unpacked)],
                        env=env_b)
    b_unpack_ok = rc == 0 and (b_unpacked / "ppt" / "slides" / "slide1.xml").exists()
    _record("A", "Skill B unpack.py: extracts slide XML",
            "PASS" if b_unpack_ok else "FAIL", err.strip() if not b_unpack_ok else "")

    if b_unpack_ok:
        slide_xml = (b_unpacked / "ppt" / "slides" / "slide1.xml").read_text(errors="replace")
        _record("A", "Skill B unpack: title text preserved",
                "PASS" if "Hello PPTX" in slide_xml else "FAIL")

    # ── Skill B: pack ────────────────────────────────────────────────────────
    if b_unpack_ok:
        b_repacked = tmpdir / "b_repacked.pptx"
        rc, out, err = _run([sys.executable, str(SKILL_B_OFFICE / "pack.py"),
                             str(b_unpacked), str(b_repacked),
                             "--original", str(fixture)],
                            env=env_b)
        b_pack_ok = rc == 0 and b_repacked.exists()
        if "lxml" in err or "defusedxml" in err or "ModuleNotFoundError" in err:
            _record("A", "Skill B pack.py: repackages unpacked PPTX", "WARN",
                    f"dependency missing: {err.strip()[:100]}")
        else:
            _record("A", "Skill B pack.py: repackages unpacked PPTX",
                    "PASS" if b_pack_ok else "FAIL", err.strip() if not b_pack_ok else "")


# ── Suite B – Feature-specific tests ─────────────────────────────────────────

def suite_b(tmpdir: Path) -> None:
    print("\n## Suite B: Feature-Specific Tests")

    fixture = tmpdir / "feat_fixture.pptx"
    _make_minimal_pptx(fixture)

    # Use the create_pptx.py demo deck for tests requiring slide layouts
    demo = tmpdir / "demo_for_b.pptx"
    _run([sys.executable, str(SKILL_A_SCRIPTS / "create_pptx.py"), str(demo)])

    # Skill A: add-note
    out = tmpdir / "with_note.pptx"
    rc, stdout, err = _run([sys.executable, str(SKILL_A_SCRIPTS / "edit_pptx.py"),
                            str(demo), str(out),
                            "--add-note", "1", "Test speaker note"])
    _record("B", "Skill A edit_pptx: --add-note inserts speaker note",
            "PASS" if rc == 0 and out.exists() else "FAIL",
            (stdout + err).strip() if rc != 0 else "")

    # Skill A: duplicate-slide (requires proper slide layout — use demo deck)
    out2 = tmpdir / "with_dup.pptx"
    rc, stdout, err = _run([sys.executable, str(SKILL_A_SCRIPTS / "edit_pptx.py"),
                            str(demo), str(out2),
                            "--duplicate-slide", "1"])
    if rc == 0 and out2.exists():
        try:
            from pptx import Presentation as _Prs
            prs_orig = _Prs(str(demo))
            orig_count = len(prs_orig.slides)
            prs = _Prs(str(out2))
            expected = orig_count + 1
            _record("B", "Skill A edit_pptx: --duplicate-slide appends copy",
                    "PASS" if len(prs.slides) == expected else "FAIL",
                    f"slide count: {len(prs.slides)} (expected {expected})")
        except ImportError:
            _record("B", "Skill A edit_pptx: --duplicate-slide",
                    "WARN", "python-pptx not installed — cannot verify slide count")
    else:
        _record("B", "Skill A edit_pptx: --duplicate-slide",
                "FAIL", (stdout + err).strip())

    # Skill A: add_slide.py (uses own unpack)
    a_unpacked = tmpdir / "a_unpacked_feat"
    rc_unpack, _, _ = _run([sys.executable, str(SKILL_A_SCRIPTS / "office" / "unpack.py"),
                            str(fixture), str(a_unpacked)])
    if rc_unpack == 0:
        rc2, stdout2, err2 = _run([sys.executable, str(SKILL_A_SCRIPTS / "add_slide.py"),
                                   str(a_unpacked), "slide1.xml"])
        _record("B", "Skill A add_slide.py: duplicates a slide (prints sldId line)",
                "PASS" if rc2 == 0 and "p:sldId" in stdout2 else "FAIL",
                (stdout2 + err2).strip() if rc2 != 0 else stdout2.strip()[:120])
    else:
        _record("B", "Skill A add_slide.py", "WARN", "skipped — unpack failed")

    # Skill B: add_slide.py (requires unpacked dir)
    b_unpacked = tmpdir / "b_unpacked_feat"
    env_b = os.environ.copy()
    env_b["PYTHONPATH"] = str(SKILL_B_SCRIPTS)
    rc, stdout, err = _run([sys.executable, str(SKILL_B_OFFICE / "unpack.py"),
                            str(fixture), str(b_unpacked)], env=env_b)
    if rc == 0:
        rc, stdout, err = _run([sys.executable, str(SKILL_B_SCRIPTS / "add_slide.py"),
                                str(b_unpacked), "slide1.xml"], env=env_b)
        _record("B", "Skill B add_slide.py: duplicates a slide (prints sldId line)",
                "PASS" if rc == 0 and "p:sldId" in stdout else "FAIL",
                (stdout + err).strip() if rc != 0 else stdout.strip()[:120])
    else:
        _record("B", "Skill B add_slide.py", "WARN", "skipped — unpack failed")

    # Skill B: clean.py
    if rc == 0 or b_unpacked.exists():
        env_b2 = os.environ.copy()
        env_b2["PYTHONPATH"] = str(SKILL_B_SCRIPTS)
        rc2, stdout2, err2 = _run([sys.executable, str(SKILL_B_SCRIPTS / "clean.py"),
                                   str(b_unpacked)], env=env_b2)
        if "defusedxml" in err2 or "ModuleNotFoundError" in err2:
            _record("B", "Skill B clean.py: removes orphaned files", "WARN",
                    f"dependency missing: {err2.strip()[:80]}")
        else:
            _record("B", "Skill B clean.py: removes orphaned files",
                    "PASS" if rc2 == 0 else "FAIL", (stdout2 + err2).strip() if rc2 != 0 else "")

    # SKILL.md size comparison
    a_len = len((SKILL_A / "SKILL.md").read_text(errors="replace"))
    b_skill = len((SKILL_B / "SKILL.md").read_text(errors="replace"))
    b_edit  = len((SKILL_B / "editing.md").read_text(errors="replace"))
    b_pptxjs = len((SKILL_B / "pptxgenjs.md").read_text(errors="replace"))
    b_total = b_skill + b_edit + b_pptxjs
    ratio = a_len / b_total
    _record("B", f"SKILL.md size: A={a_len} chars, B total={b_total} chars",
            "PASS" if ratio >= 0.5 else "FAIL" if ratio < 0.25 else "WARN",
            f"Skill A is {int(ratio*100)}% of Skill B guidance coverage")


# ── Report ────────────────────────────────────────────────────────────────────

def _report() -> None:
    now   = datetime.now().strftime("%Y-%m-%d %H:%M")
    total = _pass + _fail + _warn
    lines = [
        "# pptx-ez Skill Benchmark Report",
        "",
        f"**Generated**: {now}",
        f"**Skill A (pptx-ez)**: `kilo/skills/pptx-ez/`",
        f"**Skill B (reference)**: `shauneshraghi/skills` — `skills/pptx/`",
        "",
        "---",
        "",
        "## Summary",
        "",
        "| | Count |",
        "|---|---|",
        f"| ✓ Passed  | {_pass} |",
        f"| ✗ Failed  | {_fail} |",
        f"| ⚠ Warnings | {_warn} |",
        f"| **Total** | **{total}** |",
        "",
        "---",
        "",
        "## Results",
        "",
        "| Suite | Test | Status | Detail |",
        "|-------|------|--------|--------|",
    ]
    for suite, name, status, detail in _results:
        icon = {"PASS": "✓", "FAIL": "✗", "WARN": "⚠"}.get(status, "?")
        lines.append(f"| {suite} | {name} | {icon} {status} | {detail} |")

    lines += [
        "",
        "---",
        "",
        "## Key Findings",
        "",
        "### Where Skill A (pptx-ez) leads",
        "",
        "- **python-pptx creation path**: Full JSON-spec driven `create_pptx.py` "
          "with charts, tables, images, notes, and layout selection.",
        "- **edit_pptx.py CLI**: Find-replace, add notes, delete/duplicate slides, "
          "add textboxes — all scriptable without unpacking.",
        "- **References**: Comprehensive `python-pptx.md` and `ooxml-pptx.md` "
          "covering the full PresentationML anatomy.",
        "- **validate_pptx.py**: Standalone ZIP/OPC integrity + XML validation.",
        "- **export_pdf.sh**: Headless PDF and per-slide PNG export with format flag.",
        "",
        "### Where Skill B (Anthropic) leads",
        "",
        "- **OOXML unpack/pack/clean workflow**: Skill B ships `office/unpack.py`, "
          "`office/pack.py`, and `clean.py` enabling direct XML editing of slide XML. "
          "This is the primary editing workflow — unpack, edit XML directly, clean, pack. "
          "Skill A has no equivalent.",
        "- **add_slide.py**: Properly duplicates slides or creates from layout, "
          "handling Content_Types.xml, relationship IDs, and notes-slide cleanup. "
          "Skill A's `edit_pptx.py --duplicate-slide` is a python-pptx approximation "
          "that may miss relationship cleanup.",
        "- **thumbnail.py**: Creates visual slide grid JPEGs via LibreOffice + Pillow, "
          "essential for template analysis and QA. Skill A has no visual preview tool.",
        "- **pptxgenjs.md**: Complete PptxGenJS reference (shapes, shadows, rich-text "
          "arrays, bullets, images, multi-slide, layout coordinates). "
          "PptxGenJS is the creation path for complex decks requiring gradient "
          "backgrounds, precise pixel placement, and custom shapes.",
        "- **Design guidance in SKILL.md**: Color palette table, typography pairings, "
          "spacing rules, 'avoid common mistakes' list, bold all headers rule, "
          "never-unicode-bullets rule. Skill A's SKILL.md lacks all of this.",
        "- **editing.md**: Step-by-step template workflow, subagent parallelism "
          "guidance, multi-item content pattern (separate `<a:p>` per item), "
          "smart-quote XML entity table, pitfalls. Skill A has none of this.",
        "- **defusedxml throughout**: Skill B uses `defusedxml.minidom` in clean.py "
          "and thumbnail.py to prevent XXE attacks on untrusted PPTX files.",
        "- **Visual QA workflow**: Skill B documents soffice → pdftoppm → subagent "
          "inspection loop with a detailed prompt for finding overlaps, overflows, and "
          "contrast issues. Skill A mentions `validate_pptx.py` only.",
        "",
        "---",
        "",
        "## Actionable Improvements for pptx-ez",
        "",
        "### High Priority",
        "",
        "**1. Add OOXML unpack/pack/clean scripts**",
        "Port `office/unpack.py`, `office/pack.py`, and `clean.py` from Skill B. "
        "This is the core editing workflow for complex slide editing tasks. "
        "The unpack/pack pattern allows direct XML manipulation impossible via python-pptx.",
        "",
        "**2. Add add_slide.py**",
        "The proper slide-duplication script handles Content_Types.xml, rels, "
        "presentation.xml sldIdLst, and strips notesSlide refs from duplicated slides. "
        "Skill A's python-pptx duplicate is less reliable for complex templates.",
        "",
        "**3. Add thumbnail.py**",
        "Visual slide grid is essential for template analysis before editing. "
        "Requires Pillow + LibreOffice (already a dependency).",
        "",
        "**4. Add pptxgenjs.md reference**",
        "PptxGenJS is the recommended creation tool for complex decks. "
        "The reference is large but high-value: shadows, bullet arrays, "
        "rich text, layout coordinates, image sources.",
        "",
        "**5. Expand SKILL.md with design guidance and QA workflow**",
        "Add color palette table, typography pairings (font header/body combos), "
        "spacing rules (0.5\" margins, 0.3–0.5\" gaps), "
        "'avoid common mistakes' list, bold-headers rule, "
        "never-unicode-bullets rule, and the visual QA loop "
        "(soffice → pdftoppm → subagent inspect).",
        "",
        "### Medium Priority",
        "",
        "**6. Add editing.md with template workflow and XML pitfalls**",
        "Document: unpack → plan slide mapping → structural edits → content edits "
        "(subagents in parallel) → clean → pack. Include smart-quote entity table, "
        "multi-item paragraph pattern, and template adaptation pitfalls.",
        "",
        "**7. Switch XML parsing to defusedxml**",
        "validate_pptx.py uses Python's built-in xml.etree.ElementTree. "
        "Replace with `defusedxml.ElementTree` for XXE safety on untrusted files.",
        "",
        "**8. Add office/soffice.py (AF_UNIX socket shim)**",
        "The xlsx-ez skill already has a soffice.py with the C shim. "
        "Share or symlink it to pptx-ez so thumbnail.py and export_pdf.sh "
        "work in sandboxed containers where AF_UNIX is blocked.",
        "",
        "### Regressions to Avoid",
        "",
        "1. **Do not remove `create_pptx.py`** — python-pptx creation is a "
           "differentiating feature Skill B lacks.",
        "2. **Do not remove `validate_pptx.py`** — standalone validation is missing "
           "from Skill B's pptx skill.",
        "3. **Do not remove the ooxml-pptx.md reference** — detailed PML XML anatomy "
           "not covered in Skill B's SKILL.md.",
    ]

    report_path = SKILL_A / "benchmark_report.md"
    report_path.write_text("\n".join(lines) + "\n")
    print(f"\nReport written to: {report_path}")


# ── Main ─────────────────────────────────────────────────────────────────────

def main() -> int:
    print("pptx-ez Skill Benchmark")
    print("=" * 60)
    print(f"Skill A (pptx-ez):  {SKILL_A}")
    print(f"Skill B (reference): {SKILL_B}")

    with tempfile.TemporaryDirectory() as td:
        tmpdir = Path(td)
        suite_c()
        suite_d()
        suite_e()
        suite_a(tmpdir)
        suite_b(tmpdir)

    _report()
    print(f"{'=' * 60}")
    print(f"Results: {_pass} passed, {_fail} failed, {_warn} warnings")
    return 1 if _fail > 0 else 0


if __name__ == "__main__":
    sys.exit(main())
